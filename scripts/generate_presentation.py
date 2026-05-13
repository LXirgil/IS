from pptx import Presentation
from pptx.util import Pt
from datetime import datetime

prs = Presentation()

# タイトルスライド
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "アプリ企画と機能概要"
slide.placeholders[1].text = f"生成日: {datetime.now().strftime('%Y-%m-%d')}"

# ヘルパー: 箇条書きスライド
def add_bullets(title, bullets):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        p = body.add_paragraph()
        p.text = b
        p.level = 0

# スライド内容（リポジトリの構成を踏まえた推定）
add_bullets("アプリ概要", [
    "目的: アイデア・メモ管理とクーポン閲覧を両立するモバイル/デスクトップアプリ",
    "対象: アイデアを素早く記録したい個人、クーポン情報を一覧で管理したいユーザー",
    "価値: 思考の整理とお得情報の管理を一つのアプリで提供"])

add_bullets("コア機能", [
    "アイデア/ノート作成・編集・削除（Compose）",
    "一覧表示（ListPage）と個別詳細表示（DetailPage）",
    "クーポン一覧・クーポン詳細（CouponList / CouponDetail）",
    "検索・フィルタリング、並び替え",
    "設定画面（Settings）、データのエクスポート/共有",
    "AI支援: 要約・タグ付け・提案（ai_service.dart を想定）"])

add_bullets("画面構成とユーザーフロー", [
    "メイン -> 一覧 -> 詳細 -> 編集/作成 -> 保存",
    "設定/同期画面、検索バー、フィルタモーダル",
    "プラットフォーム別: Android/iOS/Web/Desktop 対応（Flutter）"])

add_bullets("技術スタック", [
    "UI: Flutter + Dart",
    "ローカル保存: SharedPreferences / SQLite 等",
    "AI: ローカル/外部API（ai_service.dart）",
    "対応: Android, iOS, Web, Windows, macOS, Linux"])

add_bullets("今後の開発ステップ", [
    "MVP リリース: コア機能の安定化",
    "ユーザーテストとフィードバック反映",
    "認証・プライバシー、バックアップ/同期の検討",
    "CI/CD とリリースパイプライン設定"])

prs.save('App_Plan.pptx')
print('Saved App_Plan.pptx')
