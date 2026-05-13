import re
import shutil
import os
from datetime import datetime
from docx import Document
from pptx import Presentation

# Paths
DOCX_PATH = r"c:\Users\14092\IS\IdeaStream_企画書.docx"
PPTX_PATH = r"c:\Users\14092\IS\IdeaStream.pptx"

# Backup existing PPTX
backup_path = PPTX_PATH.replace('.pptx', f'_backup_{datetime.now().strftime("%Y%m%d%H%M%S")}.pptx')
shutil.copy2(PPTX_PATH, backup_path)
print(f'Backup created: {backup_path}')

# Read docx
doc = Document(DOCX_PATH)
paras = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
text = '\n'.join(paras)

# Simple sentence split for Japanese and punctuation
sents = re.split(r'(?<=[。．.!?！？])\s*', text)
sents = [s.strip() for s in sents if s.strip()]

# Select up to 6 concise bullets
bullets = []
for s in sents:
    # skip very short fragments
    if len(s) < 6:
        continue
    bullets.append(s)
    if len(bullets) >= 6:
        break

# Fallback to paragraphs if no sentences
if not bullets:
    bullets = paras[:6]

# Make bullets shorter if too long
def shorten(s, limit=120):
    return s if len(s) <= limit else s[:limit].rstrip() + '…'

bullets = [shorten(b) for b in bullets]

# Open and update presentation
prs = Presentation(PPTX_PATH)

# Add a title slide for the doc summary
try:
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = 'IdeaStream — 企画書概要'
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = f'抽出日: {datetime.now().strftime("%Y-%m-%d")}'
except Exception:
    # fallback: create a blank slide and add a textbox
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = 'IdeaStream — 企画書概要'

# Add bullets slide
bullet_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
slide = prs.slides.add_slide(bullet_layout)
slide.shapes.title.text = '企画書の要点'

textbox = slide.shapes.placeholders[1].text_frame
textbox.clear()
for b in bullets:
    p = textbox.add_paragraph()
    p.text = b
    p.level = 0

# Save updated presentation (try overwrite; fall back to alternate name on permission error)
try:
    prs.save(PPTX_PATH)
    print(f'Updated presentation saved: {PPTX_PATH}')
except PermissionError:
    alt_path = PPTX_PATH.replace('.pptx', '_updated.pptx')
    prs.save(alt_path)
    print(f'Permission denied when writing {PPTX_PATH}; saved as {alt_path}')

print('Bullets added:')
for i, b in enumerate(bullets, 1):
    print(f'{i}. {b}')
